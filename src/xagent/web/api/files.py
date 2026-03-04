import asyncio
import logging
from pathlib import Path
from typing import Any, Dict, Optional, Tuple

from fastapi import APIRouter, Depends, File, Form, HTTPException, Query, UploadFile
from fastapi.responses import FileResponse, HTMLResponse, StreamingResponse
from pptx import Presentation
from sqlalchemy.orm import Session

from ...core.tools.adapters.vibe.file_tool import read_file
from ..auth_dependencies import get_current_user
from ..config import MAX_FILE_SIZE, UPLOADS_DIR, get_upload_path, is_allowed_file
from ..models.database import get_db
from ..models.uploaded_file import UploadedFile
from ..models.user import User
from .legacy_file import (
    infer_user_id_from_legacy_path,
    is_valid_uuid,
    resolve_legacy_file_path,
    resolve_legacy_file_path_cross_user,
)

logger = logging.getLogger(__name__)

file_router = APIRouter(prefix="/api/files", tags=["files"])


@file_router.post("/upload")
async def upload_file(
    file: UploadFile = File(...),
    task_type: str = Form(...),
    message: str = Form(""),
    task_id: str = Form(None),
    folder: str = Form(None),
    user: User = Depends(get_current_user),
) -> Dict:
    """
    Upload a single file and optionally create a task (backward compatibility)

    Args:
        file: Uploaded file
        task_type: Type of task (e.g., "general")
        message: Optional message to include with the task

    Returns:
        Upload result with file info and optional task ID
    """
    try:
        uploaded_files = []

        # Process the file
        if not file.filename or not file.filename.strip():
            raise HTTPException(status_code=422, detail="No filename provided")

        # Check file extension
        if not is_allowed_file(file.filename, task_type):
            raise HTTPException(
                status_code=500,
                detail=f"File type {Path(file.filename).suffix.lower()} not supported for task type {task_type}",
            )

        # Check file size
        content = await file.read()
        if len(content) > MAX_FILE_SIZE:
            raise HTTPException(
                status_code=500,
                detail=f"File size exceeds maximum limit of {MAX_FILE_SIZE // (1024 * 1024)}MB",
            )

        # Get upload path with user isolation
        file_path = get_upload_path(file.filename, task_id, folder, int(user.id))

        # Save uploaded file
        with open(file_path, "wb") as buffer:
            buffer.write(content)

        logger.info(f"File uploaded: {file.filename} -> {file_path} (user: {user.id})")

        # Read file content for processing
        try:
            file_content = read_file(str(file_path))

            uploaded_files.append(
                {
                    "filename": file.filename,
                    "file_path": str(file_path),
                    "file_url": get_file_url(
                        file.filename, task_id, folder, int(user.id)
                    ),
                    "file_size": len(content),
                    "content_preview": file_content[:500] + "..."
                    if len(file_content) > 500
                    else file_content,
                }
            )

        except (ValueError, KeyError, TypeError) as e:
            # Data format error
            logger.error(f"Data format error processing file {file.filename}: {e}")
            uploaded_files.append(
                {
                    "filename": file.filename,
                    "file_path": str(file_path),
                    "file_url": get_file_url(
                        file.filename, task_id, folder, int(user.id)
                    ),
                    "file_size": len(content),
                    "error": f"Data format error: {str(e)}",
                }
            )
        except (PermissionError, OSError) as e:
            # File system permission error
            logger.error(f"File system error processing file {file.filename}: {e}")
            uploaded_files.append(
                {
                    "filename": file.filename,
                    "file_path": str(file_path),
                    "file_url": get_file_url(
                        file.filename, task_id, folder, int(user.id)
                    ),
                    "file_size": len(content),
                    "error": f"File system error: {str(e)}",
                }
            )
        except Exception as e:
            # Other errors, re-raise
            logger.error(f"Unexpected error processing file {file.filename}: {e}")
            raise

        # Determine overall success
        all_successful = all("error" not in file_info for file_info in uploaded_files)

        # Return single file format
        file_info = uploaded_files[0]
        return {
            "success": all_successful,
            "filename": file_info["filename"],
            "file_path": file_info["file_path"],
            "file_url": file_info["file_url"],
            "file_size": file_info["file_size"],
            "task_type": task_type,
            "content_preview": file_info.get("content_preview", ""),
            "error": file_info.get("error"),
            "message": f"Successfully uploaded {file_info['filename']}"
            if all_successful
            else f"Failed to process {file_info['filename']}",
        }

    except HTTPException:
        # Re-raise HTTP exceptions (like 422 validation errors)
        raise
    except (ValueError, KeyError, TypeError) as e:
        # Data format error
        logger.error(f"Data format error in file upload: {e}")
        raise HTTPException(status_code=400, detail=f"Data format error: {str(e)}")
    except (PermissionError, OSError) as e:
        # File system permission error
        logger.error(f"File system error in file upload: {e}")
        raise HTTPException(status_code=403, detail=f"File system error: {str(e)}")
    except Exception as e:
        # Other errors, re-raise
        logger.error(f"Unexpected error in file upload: {e}")
        raise HTTPException(status_code=500, detail=f"Internal server error: {str(e)}")


@file_router.post("/upload-multiple")
async def upload_multiple_files(
    files: list[UploadFile] = File(...),
    task_type: str = Form(...),
    message: str = Form(""),
    task_id: str = Form(None),
    folder: str = Form(None),
    user: User = Depends(get_current_user),
) -> Dict:
    """
    Upload multiple files and optionally create a task (new format)

    Args:
        files: List of uploaded files
        task_type: Type of task (e.g., "general")
        message: Optional message to include with the task
        task_id: Optional task ID to organize files in task-specific folder
        folder: Optional folder name within task directory (e.g., "input", "output")

    Returns:
        Upload result with file info and optional task ID
    """
    try:
        uploaded_files = []

        # Process each file
        for file in files:
            # Validate file type
            if not file.filename or not file.filename.strip():
                raise HTTPException(status_code=422, detail="No filename provided")

            # Check file extension
            if not is_allowed_file(file.filename, task_type):
                raise HTTPException(
                    status_code=500,
                    detail=f"File type {Path(file.filename).suffix.lower()} not supported for task type {task_type}",
                )

            # Check file size
            content = await file.read()
            if len(content) > MAX_FILE_SIZE:
                raise HTTPException(
                    status_code=500,
                    detail=f"File size exceeds maximum limit of {MAX_FILE_SIZE // (1024 * 1024)}MB",
                )

            # Get upload path with user isolation
            file_path = get_upload_path(file.filename, task_id, folder, int(user.id))

            # Save uploaded file
            with open(file_path, "wb") as buffer:
                buffer.write(content)

            logger.info(
                f"File uploaded: {file.filename} -> {file_path} (user: {user.id})"
            )

            # Read file content for processing
            try:
                file_content = read_file(str(file_path))

                uploaded_files.append(
                    {
                        "filename": file.filename,
                        "file_path": str(file_path),
                        "file_url": get_file_url(
                            file.filename, task_id, folder, int(user.id)
                        ),
                        "file_size": len(content),
                        "content_preview": file_content[:500] + "..."
                        if len(file_content) > 500
                        else file_content,
                    }
                )

            except (ValueError, KeyError, TypeError) as e:
                # Data format error
                logger.error(f"Data format error processing file {file.filename}: {e}")
                uploaded_files.append(
                    {
                        "filename": file.filename,
                        "file_path": str(file_path),
                        "file_url": get_file_url(
                            file.filename, task_id, folder, int(user.id)
                        ),
                        "file_size": len(content),
                        "error": f"Data format error: {str(e)}",
                    }
                )
            except (PermissionError, OSError) as e:
                # File system permission error
                logger.error(f"File system error processing file {file.filename}: {e}")
                uploaded_files.append(
                    {
                        "filename": file.filename,
                        "file_path": str(file_path),
                        "file_url": get_file_url(
                            file.filename, task_id, folder, int(user.id)
                        ),
                        "file_size": len(content),
                        "error": f"File system error: {str(e)}",
                    }
                )
            except Exception as e:
                # Other errors, re-raise
                logger.error(f"Unexpected error processing file {file.filename}: {e}")
                raise

        # Determine overall success
        all_successful = all("error" not in file_info for file_info in uploaded_files)

        # Return new format for multiple files
        return {
            "success": all_successful,
            "files": uploaded_files,
            "total_files": len(uploaded_files),
            "task_type": task_type,
            "message": f"Successfully uploaded {len(uploaded_files)} files"
            if all_successful
            else "Some files had processing errors",
        }

    except HTTPException:
        # Re-raise HTTP exceptions (like 422 validation errors)
        raise
    except (ValueError, KeyError, TypeError) as e:
        # Data format error
        logger.error(f"Data format error in multiple file upload: {e}")
        raise HTTPException(status_code=400, detail=f"Data format error: {str(e)}")
    except (PermissionError, OSError) as e:
        # File system permission error
        logger.error(f"File system error in multiple file upload: {e}")
        raise HTTPException(status_code=403, detail=f"File system error: {str(e)}")
    except Exception as e:
        # Other errors, re-raise
        logger.error(f"Unexpected error in multiple file upload: {e}")
        raise HTTPException(status_code=500, detail=f"Internal server error: {str(e)}")


@file_router.get("/list")
async def list_files(user: User = Depends(get_current_user)) -> Dict:
    """List user's uploaded files"""
    try:
        # Get user-specific directory, or all directories for admin
        if user.is_admin:
            # Admin can see all files - scan all user directories
            scan_dirs = [
                d
                for d in UPLOADS_DIR.iterdir()
                if d.is_dir() and d.name.startswith("user_")
            ]
        else:
            # Regular users can only see their own files
            user_dir = UPLOADS_DIR / f"user_{user.id}"
            scan_dirs = [user_dir] if user_dir.exists() else []

        files = []

        def scan_directory(
            directory: Path, relative_path: str = "", user_prefix: str = ""
        ) -> None:
            """Recursively scan directory for files"""
            for item in directory.iterdir():
                if item.name.startswith("."):
                    continue

                if item.is_file():
                    file_info = get_file_info(str(item))
                    if file_info:
                        # Add relative path for display
                        if relative_path:
                            file_info["relative_path"] = f"{relative_path}/{item.name}"
                        else:
                            file_info["relative_path"] = item.name

                        # Parse relative_path to extract task_id and folder
                        rel_path = file_info.get("relative_path", item.name)
                        path_parts = rel_path.split("/")

                        if len(path_parts) >= 3 and path_parts[0].startswith(
                            "web_task_"
                        ):
                            # Format: web_task_X/folder/filename
                            task_id = path_parts[0].replace("web_task_", "")
                            folder = "/".join(path_parts[1:-1])
                            filename = path_parts[-1]

                            # Extract user_id from user_prefix for admin
                            file_user_id = (
                                int(user_prefix.replace("user_", ""))
                                if user_prefix
                                else int(user.id)
                            )
                            file_info["file_url"] = get_file_url(
                                filename,
                                task_id=task_id,
                                folder=folder,
                                user_id=file_user_id,
                            )
                        else:
                            # Fallback to simple filename
                            file_user_id = (
                                int(user_prefix.replace("user_", ""))
                                if user_prefix
                                else int(user.id)
                            )
                            file_info["file_url"] = get_file_url(
                                file_info["filename"], user_id=file_user_id
                            )

                        # Add user info for admin to identify file ownership
                        if user.is_admin and user_prefix:
                            file_info["user_id"] = int(user_prefix.replace("user_", ""))

                        files.append(file_info)
                elif item.is_dir():
                    # Recursively scan subdirectories
                    new_relative_path = (
                        f"{relative_path}/{item.name}" if relative_path else item.name
                    )
                    scan_directory(item, new_relative_path, user_prefix)

        # Scan all user directories
        for scan_dir in scan_dirs:
            user_prefix = scan_dir.name if scan_dir.name.startswith("user_") else ""
            scan_directory(scan_dir, "", user_prefix)

        return {"files": files, "total_count": len(files)}

    except (PermissionError, OSError) as e:
        # File system permission error
        logger.error(f"File system error listing files: {e}")
        raise HTTPException(status_code=403, detail=f"File system error: {str(e)}")
    except (ValueError, KeyError, TypeError) as e:
        # Data format error
        logger.error(f"Data format error listing files: {e}")
        raise HTTPException(status_code=400, detail=f"Data format error: {str(e)}")
    except Exception as e:
        # Other errors, re-raise
        logger.error(f"Unexpected error listing files: {e}")
        raise HTTPException(status_code=500, detail=f"Internal server error: {str(e)}")


@file_router.get("/download/{file_path:path}", response_model=None)
async def download_file(
    file_path: str,
    user: User = Depends(get_current_user),
    db: Session = Depends(get_db),
) -> Any:
    """Download uploaded file"""
    try:
        logger.info(f"Download request for file_path: {file_path}")

        # Check if this is a web_task file and handle admin access
        target_user_dir = None
        if file_path.startswith("web_task_"):
            # Extract task ID from path like "web_task_78/output/file.html"
            try:
                task_id = int(file_path.split("_")[2].split("/")[0])
                from ..models.task import Task

                task = db.query(Task).filter(Task.id == task_id).first()

                if task:
                    # Admin can access any task, regular users can only access their own tasks
                    if not user.is_admin and task.user_id != int(user.id):
                        logger.warning(
                            f"User {user.id} attempted to access task {task_id} belonging to user {task.user_id}"
                        )
                        raise HTTPException(status_code=403, detail="Access denied")

                    # Use the task owner's directory for file access
                    target_user_dir = UPLOADS_DIR / f"user_{task.user_id}"
                    logger.info(
                        f"Accessing task {task_id} files from user directory: user_{task.user_id}"
                    )
                else:
                    logger.warning(f"Task {task_id} not found in database")
                    raise HTTPException(status_code=404, detail="Task not found")
            except (ValueError, IndexError) as e:
                logger.warning(f"Invalid web_task path format: {file_path}, error: {e}")
                # Fall back to normal processing if path format is invalid
                pass

        # Get user-specific directory
        # For admin accessing other users' task files, use the target user's directory
        if target_user_dir:
            user_dir = target_user_dir
        else:
            user_dir = UPLOADS_DIR / f"user_{user.id}"

        # Define recursive search function once
        def find_file_recursively(directory: Path, filename: str) -> Path | None:
            """Recursively find file in directory"""
            for item in directory.iterdir():
                if item.is_file() and item.name == filename:
                    return item
                elif item.is_dir():
                    result = find_file_recursively(item, filename)
                    if result:
                        return result
            return None

        if "/" not in file_path:
            # Simple filename - search in user directory first, then fallback to global
            found_path = None

            # Search in user directory first
            if user_dir.exists():
                found_path = find_file_recursively(user_dir, file_path)
                if found_path:
                    full_path = found_path
                    logger.info(f"Found file in user directory: {full_path}")

            # If not found in user directory, search in global directory for legacy files
            if not found_path:
                if UPLOADS_DIR.exists():
                    found_path = find_file_recursively(UPLOADS_DIR, file_path)
                    if found_path:
                        full_path = found_path
                        logger.info(
                            f"Found file in global directory (legacy): {full_path}"
                        )

            if not found_path:
                logger.warning(f"File not found for download: {file_path}")
                raise HTTPException(status_code=404, detail="File not found")
        else:
            # Relative path - check if it starts with user_id prefix
            if file_path.startswith(f"user_{user.id}/"):
                # Path already includes user directory, use it relative to UPLOADS_DIR
                full_path = UPLOADS_DIR / file_path
            elif target_user_dir:
                # For admin accessing other users' task files, path is relative to target user directory
                full_path = target_user_dir / file_path
            else:
                # Path is relative to user directory
                full_path = user_dir / file_path

            if not full_path.exists():
                # Fallback: search recursively for the filename with path preference
                target_filename = file_path.split("/")[-1]

                # Enhanced search that prefers matching path structure
                def find_file_with_path_preference(
                    directory: Path, filename: str, original_path: str
                ) -> Path | None:
                    """Find file recursively, preferring paths that match the original structure"""
                    best_match = None
                    best_score = 0

                    def score_match(candidate_path: Path, target_path: str) -> int:
                        """Score how well a candidate path matches the target structure"""
                        # Get relative path from user directory
                        try:
                            relative_path = candidate_path.relative_to(user_dir)
                        except ValueError:
                            # Fallback to using absolute path parts
                            relative_path = candidate_path

                        candidate_parts = relative_path.parts
                        target_parts = target_path.split("/")

                        score = 0
                        # Check if path parts match
                        for i, target_part in enumerate(
                            target_parts[:-1]
                        ):  # Exclude filename
                            if (
                                i < len(candidate_parts) - 1
                            ):  # Exclude filename from candidate
                                candidate_part = candidate_parts[i]
                                if target_part == candidate_part:
                                    score += 10
                                elif (
                                    target_part.replace("web_task_", "task_")
                                    == candidate_part
                                ):
                                    score += (
                                        5  # Partial match for task prefix differences
                                    )
                                elif (
                                    candidate_part.replace("web_task_", "task_")
                                    == target_part
                                ):
                                    score += 5  # Reverse match
                        return score

                    for item in directory.iterdir():
                        if item.is_file() and item.name == filename:
                            # Score this match
                            current_score = score_match(item, original_path)
                            if current_score > best_score:
                                best_match = item
                                best_score = current_score
                                logger.debug(
                                    f"Found better match with score {current_score}: {item}"
                                )
                        elif item.is_dir():
                            result = find_file_with_path_preference(
                                item, filename, original_path
                            )
                            if result:
                                current_score = score_match(result, original_path)
                                if current_score > best_score:
                                    best_match = result
                                    best_score = current_score

                    return best_match

                found_path = find_file_with_path_preference(
                    user_dir, target_filename, file_path
                )
                if found_path:
                    full_path = found_path
                    logger.info(f"Found file by enhanced recursive search: {full_path}")
                else:
                    logger.warning(f"File not found for download: {file_path}")
                    raise HTTPException(status_code=404, detail="File not found")
            else:
                logger.info(f"Found file by direct path: {full_path}")

        # Security check: ensure the path is within allowed directories
        try:
            resolved_path = full_path.resolve()
            resolved_user_dir = user_dir.resolve()
            resolved_global_dir = UPLOADS_DIR.resolve()

            # Allow both user directory and global directory for legacy files
            try:
                resolved_path.relative_to(resolved_user_dir)
            except ValueError:
                try:
                    resolved_path.relative_to(resolved_global_dir)
                except ValueError:
                    logger.warning(f"Security check failed for path: {file_path}")
                    raise HTTPException(status_code=403, detail="Access denied")
        except ValueError:
            logger.warning(f"Security check failed for path: {file_path}")
            raise HTTPException(status_code=403, detail="Access denied")

        if not full_path.exists():
            logger.warning(f"File not found for download: {full_path}")
            raise HTTPException(status_code=404, detail="File not found")

        filename = full_path.name

        # Handle PPTX files - convert to PDF using LibreOffice for preview
        if filename.endswith(".pptx"):
            logger.info(f"Converting PPTX to PDF for download: {full_path}")
            import tempfile

            try:
                # Create a temporary directory for the conversion
                with tempfile.TemporaryDirectory() as temp_dir:
                    # Use LibreOffice to convert PPTX to PDF (async)
                    proc = await asyncio.create_subprocess_exec(
                        "soffice",
                        "--headless",
                        "--convert-to",
                        "pdf",
                        "--outdir",
                        temp_dir,
                        str(full_path),
                        stdout=asyncio.subprocess.PIPE,
                        stderr=asyncio.subprocess.PIPE,
                    )

                    try:
                        stdout, stderr = await asyncio.wait_for(
                            proc.communicate(), timeout=30
                        )
                    except asyncio.TimeoutError:
                        proc.kill()
                        await proc.wait()
                        logger.error(
                            f"LibreOffice conversion timed out for PPTX: {full_path}"
                        )
                        raise

                    if proc.returncode == 0:
                        # Find the generated PDF file
                        pdf_files = list(Path(temp_dir).glob("*.pdf"))
                        if pdf_files:
                            pdf_path = pdf_files[0]
                            logger.info(
                                f"Successfully converted PPTX to PDF: {pdf_path}"
                            )

                            # Read PDF into memory before temp dir is cleaned up
                            pdf_content = pdf_path.read_bytes()

                            # Return streaming response from memory
                            return StreamingResponse(
                                iter([pdf_content]),
                                media_type="application/pdf",
                                headers={
                                    "Content-Disposition": f'inline; filename="{full_path.stem}.pdf"'
                                },
                            )
                        else:
                            logger.warning(
                                "LibreOffice conversion succeeded but no PDF found"
                            )
                    else:
                        logger.warning(
                            f"LibreOffice conversion failed: {stderr.decode()}"
                        )

            except asyncio.TimeoutError:
                logger.error(f"LibreOffice conversion timed out for PPTX: {full_path}")
            except FileNotFoundError:
                logger.warning(
                    "LibreOffice (soffice) not found, returning original PPTX"
                )
            except Exception as e:
                logger.error(f"Failed to convert PPTX to PDF: {e}")

        # Determine media type
        media_type = "application/octet-stream"
        if filename.endswith((".html", ".htm")):
            media_type = "text/html"
        elif filename.endswith((".css")):
            media_type = "text/css"
        elif filename.endswith((".js")):
            media_type = "application/javascript"
        elif filename.endswith((".jpg", ".jpeg")):
            media_type = "image/jpeg"
        elif filename.endswith((".png")):
            media_type = "image/png"
        elif filename.endswith((".gif")):
            media_type = "image/gif"
        elif filename.endswith((".svg")):
            media_type = "image/svg+xml"
        elif filename.endswith((".webp")):
            media_type = "image/webp"
        elif filename.endswith((".pdf")):
            media_type = "application/pdf"

        return FileResponse(
            path=str(full_path),
            filename=filename,
            media_type=media_type,
        )

    except HTTPException:
        # Re-raise HTTP exceptions (like 404 not found)
        raise
    except (PermissionError, OSError) as e:
        # File system error
        logger.error(f"File system error downloading file: {e}")
        raise HTTPException(status_code=403, detail=f"File system error: {str(e)}")
    except (ValueError, KeyError, TypeError) as e:
        # Data format error
        logger.error(f"Data format error downloading file: {e}")
        raise HTTPException(status_code=400, detail=f"Data format error: {str(e)}")
    except Exception as e:
        # Other errors, re-raise
        logger.error(f"Unexpected error downloading file: {e}")
        raise HTTPException(status_code=500, detail=f"Internal server error: {str(e)}")


@file_router.get("/preview/{task_id:int}/{file_path:path}", response_model=None)
async def preview_file(
    task_id: int,
    file_path: str,
    db: Session = Depends(get_db),
) -> Any:
    """
    Preview a file with automatic PPTX to HTML conversion.

    This endpoint checks if the file is a PPTX and automatically converts
    it to HTML for browser preview. For other file types, returns the file as-is.
    """
    try:
        logger.info(f"Preview request for task {task_id}, file: {file_path}")

        # Get task from database
        from ..models.task import Task

        task = db.query(Task).filter(Task.id == task_id).first()

        if not task:
            logger.warning(f"Task {task_id} not found for preview")
            raise HTTPException(status_code=404, detail="Task not found")

        # Use the task owner's directory for file access
        user_dir = UPLOADS_DIR / f"user_{task.user_id}"

        # Define recursive search function
        def find_file_recursively(directory: Path, filename: str) -> Path | None:
            """Recursively find file in directory"""
            for item in directory.iterdir():
                if item.is_file() and item.name == filename:
                    return item
                elif item.is_dir():
                    result = find_file_recursively(item, filename)
                    if result:
                        return result
            return None

        # Find the file
        found_path = None

        # Search in user directory
        if user_dir.exists():
            # Try to find by exact relative path first
            potential_path = user_dir / file_path
            if potential_path.exists():
                found_path = potential_path
            else:
                # Fallback to recursive search by filename
                filename = file_path.split("/")[-1]
                found_path = find_file_recursively(user_dir, filename)

        if not found_path:
            logger.warning(f"File not found for preview: {file_path}")
            raise HTTPException(status_code=404, detail="File not found")

        # Security check: ensure the path is within the task owner's directory
        try:
            resolved_path = found_path.resolve()
            resolved_user_dir = user_dir.resolve()
            resolved_path.relative_to(resolved_user_dir)
        except ValueError:
            logger.warning(f"Security check failed for preview path: {file_path}")
            raise HTTPException(status_code=403, detail="Access denied")

        # Get filename early for PPTX processing
        filename = found_path.name

        # Handle PPTX files - convert to PDF using LibreOffice for preview
        if filename.endswith(".pptx"):
            logger.info(f"Converting PPTX to PDF for preview: {found_path}")
            import tempfile

            try:
                # Create a temporary directory for the conversion
                with tempfile.TemporaryDirectory() as temp_dir:
                    # Use LibreOffice to convert PPTX to PDF (async)
                    proc = await asyncio.create_subprocess_exec(
                        "soffice",
                        "--headless",
                        "--convert-to",
                        "pdf",
                        "--outdir",
                        temp_dir,
                        str(found_path),
                        stdout=asyncio.subprocess.PIPE,
                        stderr=asyncio.subprocess.PIPE,
                    )

                    try:
                        stdout, stderr = await asyncio.wait_for(
                            proc.communicate(), timeout=30
                        )
                    except asyncio.TimeoutError:
                        proc.kill()
                        await proc.wait()
                        logger.error(
                            f"LibreOffice conversion timed out for PPTX: {found_path}"
                        )
                        raise

                    if proc.returncode == 0:
                        # Find the generated PDF file
                        pdf_files = list(Path(temp_dir).glob("*.pdf"))
                        if pdf_files:
                            pdf_path = pdf_files[0]
                            logger.info(
                                f"Successfully converted PPTX to PDF: {pdf_path}"
                            )

                            # Read PDF into memory before temp dir is cleaned up
                            pdf_content = pdf_path.read_bytes()

                            # Return streaming response from memory
                            return StreamingResponse(
                                iter([pdf_content]),
                                media_type="application/pdf",
                                headers={
                                    "Content-Disposition": f'inline; filename="{found_path.stem}.pdf"'
                                },
                            )
                        else:
                            logger.warning(
                                "LibreOffice conversion succeeded but no PDF found"
                            )
                    else:
                        logger.warning(
                            f"LibreOffice conversion failed: {stderr.decode()}"
                        )

            except asyncio.TimeoutError:
                logger.error(f"LibreOffice conversion timed out for PPTX: {found_path}")
            except FileNotFoundError:
                logger.warning(
                    "LibreOffice (soffice) not found, falling back to text extraction"
                )
            except Exception as e:
                logger.error(f"Failed to convert PPTX to PDF: {e}")

            # Fallback: extract text using python-pptx if LibreOffice fails
            logger.info(f"Falling back to text extraction for PPTX: {found_path}")
            try:
                prs = Presentation(str(found_path))
                html_content = """
                <!DOCTYPE html>
                <html>
                <head>
                    <meta charset="UTF-8">
                    <style>
                        body { font-family: Arial, sans-serif; margin: 40px; line-height: 1.6; }
                        h1 { color: #333; border-bottom: 2px solid #4CAF50; padding-bottom: 10px; }
                        h2 { color: #555; margin-top: 30px; }
                        .slide { border: 1px solid #ddd; padding: 20px; margin: 20px 0; background: #f9f9f9; border-radius: 8px; }
                        .slide-number { color: #999; font-size: 12px; margin-top: 10px; }
                        .text-content { white-space: pre-wrap; }
                    </style>
                </head>
                <body>
                    <h1>📊 {filename}</h1>
                """

                for slide_num, slide in enumerate(prs.slides, 1):
                    slide_text = []
                    for shape in slide.shapes:
                        if hasattr(shape, "text") and shape.text:
                            slide_text.append(shape.text)

                    if slide_text:
                        html_content += f"""
                        <div class="slide">
                            <h2>Slide {slide_num}</h2>
                            <div class="text-content">
                                {"<br>".join(slide_text)}
                            </div>
                            <div class="slide-number">Slide {slide_num} of {len(prs.slides)}</div>
                        </div>
                        """

                html_content += """
                </body>
                </html>
                """

                return HTMLResponse(content=html_content)

            except Exception as e:
                logger.error(f"Failed to extract text from PPTX: {e}")
                # Fall through to default file response

        # Determine media type
        media_type = "application/octet-stream"
        if filename.endswith((".html", ".htm")):
            media_type = "text/html"
        elif filename.endswith((".css")):
            media_type = "text/css"
        elif filename.endswith((".js")):
            media_type = "application/javascript"
        elif filename.endswith((".jpg", ".jpeg")):
            media_type = "image/jpeg"
        elif filename.endswith((".png")):
            media_type = "image/png"
        elif filename.endswith((".gif")):
            media_type = "image/gif"
        elif filename.endswith((".svg")):
            media_type = "image/svg+xml"
        elif filename.endswith((".webp")):
            media_type = "image/webp"
        elif filename.endswith((".pdf")):
            media_type = "application/pdf"

        logger.info(f"Serving preview file: {found_path}")
        return FileResponse(
            path=str(found_path),
            filename=filename,
            media_type=media_type,
            headers={"Content-Disposition": "inline"},
        )

    except HTTPException:
        # Re-raise HTTP exceptions
        raise
    except (PermissionError, OSError) as e:
        logger.error(f"File system error in preview: {e}")
        raise HTTPException(status_code=403, detail=f"File system error: {str(e)}")
    except Exception as e:
        logger.error(f"Unexpected error in preview: {e}")
        raise HTTPException(status_code=500, detail=f"Internal server error: {str(e)}")


@file_router.get("/public/preview/{task_id:int}/{file_path:path}")
async def public_preview_file(
    task_id: int,
    user: User = Depends(get_current_user),
    db: Session = Depends(get_db),
) -> Dict[str, Any]:
    """
    Get all files for a specific task.

    More efficient than /api/files/list as it filters at database level.
    Only returns files that are already registered in the database.
    """
    # Query files for this task
    query = db.query(UploadedFile).filter(UploadedFile.task_id == task_id)

    # Permission check: only show user's own files unless admin
    if not _is_admin_user(user):
        query = query.filter(UploadedFile.user_id == _user_id_value(user))

    records = query.order_by(UploadedFile.created_at.desc()).all()

    files = []
    for record in records:
        path = Path(_file_storage_path_value(record))
        if not path.exists():
            # Skip files that no longer exist on disk
            continue

        record_user_id = _file_user_id_value(record)
        relative_path = _extract_relative_path(path, record_user_id)

        # Categorize by directory (input/output/temp)
        path_parts = relative_path.split("/")
        file_category = "other"
        if len(path_parts) >= 2:
            subdir = path_parts[1]  # e.g., "input", "output", "temp"
            if subdir in ["input", "output", "temp"]:
                file_category = subdir

        files.append(
            {
                "file_id": record.file_id,
                "filename": _file_name_value(record),
                "file_size": record.file_size,
                "modified_time": _to_unix_timestamp(path, record.created_at),
                "file_type": path.suffix.lower().lstrip("."),
                "relative_path": relative_path,
                "category": file_category,
                "task_id": record.task_id,
                "user_id": record_user_id,
            }
        )

    return {"files": files, "total_count": len(files), "task_id": task_id}


@file_router.get("/download/{file_id:path}", response_model=None)
async def download_file(
    file_id: str,
    user: User = Depends(get_current_user),
    db: Session = Depends(get_db),
) -> Any:
    file_record, full_path, owner_user_id = _resolve_file_path(
        db, file_id, _user_id_value(user)
    )

    # Check access permissions
    if file_record:
        _check_file_access(file_record, user)
        file_name = _file_name_value(file_record)
        media_type = _guess_media_type(file_name)
    else:
        # For legacy files without records, check ownership
        if owner_user_id != _user_id_value(user) and not _is_admin_user(user):
            raise HTTPException(status_code=403, detail="Access denied")
        file_name = full_path.name
        media_type = _guess_media_type(file_name)

    _ensure_under_uploads(full_path, owner_user_id)

    if not full_path.exists():
        raise HTTPException(status_code=404, detail="File not found")

    converted_pdf = await _try_convert_pptx_to_pdf(full_path)
    if converted_pdf is not None:
        return converted_pdf

    # For images and other viewable content, set Content-Disposition to inline
    # to allow browser to display the file instead of downloading it
    content_disposition = (
        "inline"
        if media_type.startswith(("image/", "video/", "audio/", "text/"))
        else "attachment"
    )

    return FileResponse(
        path=str(full_path),
        filename=file_name,
        media_type=media_type,
        headers={
            "Content-Disposition": f'{content_disposition}; filename="{file_name}"'
        },
    )


@file_router.get("/preview/{file_id:path}", response_model=None)
async def preview_file(
    file_id: str,
    user: User = Depends(get_current_user),
    db: Session = Depends(get_db),
) -> Any:
    file_record, full_path, owner_user_id = _resolve_file_path(
        db, file_id, _user_id_value(user)
    )

    # Check access permissions
    if file_record:
        _check_file_access(file_record, user)
        file_name = _file_name_value(file_record)
        media_type = _guess_media_type(file_name)
    else:
        # For legacy files without records, check ownership
        if owner_user_id != _user_id_value(user) and not _is_admin_user(user):
            raise HTTPException(status_code=403, detail="Access denied")
        file_name = full_path.name
        media_type = _guess_media_type(file_name)

    _ensure_under_uploads(full_path, owner_user_id)

    if not full_path.exists():
        raise HTTPException(status_code=404, detail="File not found")

    converted_pdf = await _try_convert_pptx_to_pdf(full_path)
    if converted_pdf is not None:
        return converted_pdf

    if full_path.suffix.lower() == ".pptx":
        try:
            return _pptx_fallback_html(full_path)
        except Exception:
            pass

    return FileResponse(
        path=str(full_path),
        filename=file_name,
        media_type=media_type,
        headers={"Content-Disposition": "inline"},
    )


@file_router.get("/public/preview/{file_id:path}", response_model=None)
async def public_preview_file(
    file_id: str,
    relative_path: Optional[str] = Query(default=None),
    db: Session = Depends(get_db),
) -> Any:
    # For public preview, we need to handle both file_id and legacy paths
    # Try UUID first
    file_record = None
    base_path = None
    owner_user_id = None

    if is_valid_uuid(file_id):
        file_record = (
            db.query(UploadedFile).filter(UploadedFile.file_id == file_id).first()
        )

    if file_record:
        base_path = Path(_file_storage_path_value(file_record))
        owner_user_id = _file_user_id_value(file_record)
    else:
        # Try to resolve as legacy path across all user directories
        result = resolve_legacy_file_path_cross_user(file_id)
        if result is None:
            raise HTTPException(status_code=404, detail="File not found")

        base_path, owner_user_id = result

    target_path = _resolve_public_preview_target(
        base_path,
        relative_path,
        owner_user_id,
    )

    if not target_path.exists() or not target_path.is_file():
        raise HTTPException(status_code=404, detail="File not found")

    converted_pdf = await _try_convert_pptx_to_pdf(target_path)
    if converted_pdf is not None:
        return converted_pdf

    return FileResponse(
        path=str(target_path),
        filename=target_path.name,
        media_type=_guess_media_type(target_path.name),
        headers={"Content-Disposition": "inline"},
    )


@file_router.post("/backfill")
async def backfill_files(
    user: User = Depends(get_current_user), db: Session = Depends(get_db)
) -> Dict[str, Any]:
    """
    Manually trigger file backfill to sync filesystem with database.

    This is a maintenance operation that scans the filesystem and creates
    database records for any unregistered files. Only available to admins.
    """
    if not _is_admin_user(user):
        raise HTTPException(status_code=403, detail="Admin access required")

    try:
        _backfill_uploaded_file_records(db, user)
        return {"success": True, "message": "File backfill completed successfully"}
    except Exception as e:
        db.rollback()
        raise HTTPException(status_code=500, detail=f"Backfill failed: {str(e)}") from e


@file_router.delete("/{file_id:path}")
async def delete_file(
    file_id: str,
    user: User = Depends(get_current_user),
    db: Session = Depends(get_db),
) -> Dict[str, Any]:
    file_record, file_path, owner_user_id = _resolve_file_path(
        db, file_id, _user_id_value(user)
    )

    # Check access permissions
    if file_record:
        _check_file_access(file_record, user)
        file_name = _file_name_value(file_record)
    else:
        # For legacy files without records, check ownership
        if owner_user_id != _user_id_value(user) and not _is_admin_user(user):
            raise HTTPException(status_code=403, detail="Access denied")
        file_name = file_path.name

    _ensure_under_uploads(file_path, owner_user_id)

    if file_path.exists() and file_path.is_file():
        file_path.unlink()

    # Delete database record if exists
    if file_record:
        db.delete(file_record)
        db.commit()

    return {
        "success": True,
        "message": f"File {file_name} deleted successfully",
        "file_id": file_id,
    }
